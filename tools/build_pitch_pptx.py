#!/usr/bin/env python3
"""
Builds a PowerPoint (.pptx) from the Markdown deck at docs/pitch/SmartFarm_AI_Pitch_Deck.md
Usage (Windows):
    python tools/build_pitch_pptx.py
Output:
    docs/pitch/SmartFarm_AI_Pitch_Deck.pptx
Dependencies:
    pip install python-pptx
"""
from __future__ import annotations
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except Exception as e:
    sys.stderr.write(
        "python-pptx is not installed. Install with: pip install python-pptx\n"
    )
    raise

ROOT = Path(__file__).resolve().parents[1]
MD_PATH = ROOT / "docs" / "pitch" / "SmartFarm_AI_Pitch_Deck.md"
OUT_PATH = ROOT / "docs" / "pitch" / "SmartFarm_AI_Pitch_Deck.pptx"

SEP_RE = re.compile(r"^\s*---\s*$")
TITLE_RE = re.compile(r"^\s*#{1,6}\s*(.+?)\s*$")
BULLET_RE = re.compile(r"^\s*-\s+(.*)$")
PLACEHOLDER_RE = re.compile(r"^\s*\[Placeholder:(.*?)\]\s*$", re.IGNORECASE)


def split_slides(md_text: str) -> list[str]:
    slides: list[str] = []
    buf: list[str] = []
    for line in md_text.splitlines():
        if SEP_RE.match(line):
            if buf:
                slides.append("\n".join(buf).strip())
                buf = []
        else:
            buf.append(line)
    if buf:
        slides.append("\n".join(buf).strip())
    return [s for s in slides if s]


def parse_slide(md_chunk: str) -> tuple[str, list[str]]:
    """Return (title, bullets). If no explicit bullets, treat paragraphs as bullets."""
    lines = [l for l in md_chunk.splitlines() if l.strip()]
    title = ""
    body_lines: list[str] = []
    for i, line in enumerate(lines):
        m = TITLE_RE.match(line)
        if m:
            title = m.group(1)
            body_lines = lines[i + 1 :]
            break
    if not title:
        # Fallback: first non-empty line is the title
        title = lines[0]
        body_lines = lines[1:]

    bullets: list[str] = []
    for bl in body_lines:
        bm = BULLET_RE.match(bl)
        if bm:
            bullets.append(bm.group(1).strip())
        else:
            # Treat non-bullet lines as bullet paragraphs unless it's a placeholder
            pm = PLACEHOLDER_RE.match(bl)
            if pm:
                bullets.append(f"[Placeholder]{pm.group(1).strip()}")
            else:
                # Skip slide separators or horizontal rules
                if bl.strip('- ').strip():
                    bullets.append(bl.strip())
    return title.strip(), bullets


def add_textbox(slide, text: str, left=Inches(1.0), top=Inches(5.0), width=Inches(8.0), height=Inches(1.0)):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def build_pptx(slides_md: list[str]) -> Presentation:
    prs = Presentation()
    # 16:9 default is fine; use Title and Content layout (index 1)
    layout = prs.slide_layouts[1]
    for md in slides_md:
        title, bullets = parse_slide(md)
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        # Clear default paragraph
        body.clear()
        if not bullets:
            p = body.add_paragraph()
            p.text = ""
        for b in bullets:
            if b.startswith("[Placeholder]"):
                add_textbox(slide, b.replace("[Placeholder]", "Placeholder:").strip())
                continue
            p = body.add_paragraph()
            p.text = b
            p.level = 0
    return prs


def main() -> int:
    if not MD_PATH.exists():
        sys.stderr.write(f"Deck markdown not found: {MD_PATH}\n")
        return 1
    md = MD_PATH.read_text(encoding="utf-8")
    slides_md = split_slides(md)
    if not slides_md:
        sys.stderr.write("No slides parsed from markdown. Check '---' separators.\n")
        return 1
    prs = build_pptx(slides_md)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
